from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import List, Tuple

from bs4 import BeautifulSoup
from pypdf import PdfReader


@dataclass
class ExtractedDoc:
    title: str
    text: str
    pages: List[Tuple[int, str]]  # (page_number, page_text); empty for non-paginated formats


def _clean_whitespace(s: str) -> str:
    s = re.sub(r"[ \t ]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def extract_html(content: bytes, url: str) -> ExtractedDoc:
    soup = BeautifulSoup(content, "lxml")
    for tag in soup(["script", "style", "noscript", "nav", "footer", "header", "form", "svg"]):
        tag.decompose()

    title = (soup.title.string.strip() if soup.title and soup.title.string else url)

    main = soup.find("main") or soup.find("article") or soup.body or soup
    text_parts: List[str] = []
    for el in main.find_all(["h1", "h2", "h3", "h4", "p", "li", "blockquote", "td", "th"]):
        line = el.get_text(separator=" ", strip=True)
        if line:
            if el.name in {"h1", "h2", "h3", "h4"}:
                text_parts.append(f"\n\n# {line}\n")
            else:
                text_parts.append(line)

    text = _clean_whitespace("\n".join(text_parts))
    return ExtractedDoc(title=title, text=text, pages=[])


def extract_pdf(content: bytes, url: str) -> ExtractedDoc:
    reader = PdfReader(io.BytesIO(content))
    pages: List[Tuple[int, str]] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        txt = _clean_whitespace(txt)
        if txt:
            pages.append((i, txt))

    title = url.rsplit("/", 1)[-1]
    try:
        meta_title = reader.metadata.title if reader.metadata else None
        if meta_title:
            title = str(meta_title).strip() or title
    except Exception:
        pass

    full = "\n\n".join(f"[Page {p}]\n{t}" for p, t in pages)
    return ExtractedDoc(title=title, text=full, pages=pages)


def extract_docx(content: bytes, url: str) -> ExtractedDoc:
    from docx import Document  # python-docx

    doc = Document(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    text = _clean_whitespace("\n".join(paragraphs))
    title = url.rsplit("/", 1)[-1]
    return ExtractedDoc(title=title, text=text, pages=[])


def extract_pptx(content: bytes, url: str) -> ExtractedDoc:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    slides_text: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
        slides_text.append("\n".join(lines))
    text = _clean_whitespace("\n\n".join(slides_text))
    title = url.rsplit("/", 1)[-1]
    return ExtractedDoc(title=title, text=text, pages=[])


def extract_xlsx(content: bytes, url: str) -> ExtractedDoc:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), data_only=True, read_only=True)
    parts: List[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append("\t".join(cells))
    text = _clean_whitespace("\n".join(parts))
    title = url.rsplit("/", 1)[-1]
    return ExtractedDoc(title=title, text=text, pages=[])


def extract(content: bytes, content_type: str, url: str) -> ExtractedDoc | None:
    ct = (content_type or "").lower()
    lower_url = url.lower()
    try:
        if "pdf" in ct or lower_url.endswith(".pdf"):
            return extract_pdf(content, url)
        if (
            "officedocument.wordprocessingml" in ct
            or "msword" in ct
            or lower_url.endswith(".docx")
        ):
            return extract_docx(content, url)
        if "presentationml" in ct or lower_url.endswith(".pptx"):
            return extract_pptx(content, url)
        if "spreadsheetml" in ct or lower_url.endswith(".xlsx"):
            return extract_xlsx(content, url)
        if "html" in ct or "xml" in ct or lower_url.endswith((".html", ".htm", "/")):
            return extract_html(content, url)
    except Exception as e:
        print(f"[extractor] failed to parse {url}: {e}")
        return None
    return None
